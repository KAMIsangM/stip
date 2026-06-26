"""PPT generator — outputs slide JSON for frontend carousel rendering + .pptx file + narration MP3s."""

from __future__ import annotations

import asyncio
import json
import logging
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from app.core.config import get_assets_root
from app.generator.base_generator import BaseModalGenerator
from app.provider.factory import get_tts_provider

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Design constants — Illustration style (warm, playful, hand-drawn feel)
# ---------------------------------------------------------------------------
# Primary palette — warm orange
_COLOR_PRIMARY = RGBColor(0xFF, 0x8C, 0x42)       # Warm orange (main)
_COLOR_PRIMARY_LIGHT = RGBColor(0xFF, 0xBF, 0x80)  # Light orange
_COLOR_PRIMARY_DARK = RGBColor(0xE0, 0x6A, 0x1E)  # Dark orange for gradients

# Secondary palette — fresh green
_COLOR_SECONDARY = RGBColor(0x4E, 0xCD, 0xC4)    # Fresh green (secondary)
_COLOR_SECONDARY_LIGHT = RGBColor(0x7E, 0xFF, 0xE4)  # Light mint
_COLOR_SECONDARY_DARK = RGBColor(0x2A, 0x9D, 0x8F)  # Dark teal

# Accent palette — warm yellow / pink
_COLOR_ACCENT_YELLOW = RGBColor(0xFF, 0xE6, 0x6D)  # Warm yellow
_COLOR_ACCENT_PINK = RGBColor(0xFF, 0x6B, 0x6B)    # Coral pink
_COLOR_ACCENT_PURPLE = RGBColor(0x95, 0x7F, 0xCD)  # Soft purple

# Background & text
_COLOR_BG = RGBColor(0xFF, 0xF8, 0xF0)             # Warm white background
_COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_COLOR_TEXT = RGBColor(0x3D, 0x2C, 0x2C)            # Deep brown text
_COLOR_TEXT_LIGHT = RGBColor(0x8B, 0x7E, 0x7E)      # Light brown text

# Bullet & decorative
_COLOR_BULLET = RGBColor(0xFF, 0x8C, 0x42)         # Bullet dot (warm orange)
_COLOR_DECOR_CIRCLE = RGBColor(0xFF, 0xE6, 0x6D)    # Decorative circle (yellow)

_FONT_FAMILY = "Microsoft YaHei"  # 微软雅黑
_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)

_PPT_SYSTEM = """你是一位资深教学设计专家和PPT课件设计师，擅长制作高质量、插画风格的教学幻灯片。你的任务是根据提供的章节信息和知识点，生成一套内容详实、教学逻辑严谨、视觉丰富的PPT课件。

## 核心设计原则

1. **输出必须是严格的 JSON 格式**，不可包含任何解释性文字
2. **幻灯片数量**：根据知识点数量和复杂度灵活调整，通常 5-10 张（知识点多的章节应更多）
3. **教学逻辑顺序**：引入热身 → 概念讲解 → 深入分析 → 案例/应用 → 互动/检验 → 总结回顾
4. **知识点全覆盖**：必须覆盖所有提供的知识点，不能遗漏任何重要知识点
5. **内容要有教学深度**：不能只是标题式罗列，要有具体的解释、例证、对比和延伸
6. **布局多样化**：根据教学内容选择合适的 layout，避免所有幻灯片都是同一种布局

## 布局类型说明（layout 字段）

你必须根据每张幻灯片的内容特点，从以下布局中选择最合适的一种：

### "title" — 标题页
- 用于课件的第一张幻灯片
- 包含课件的整体标题
- 不需要 bullets，但需要有吸引人的标题文字

### "content" — 标准内容页（默认）
- 用于概念讲解、要点说明等
- 包含 bullets 要点列表
- 最适合文字内容较多的页面

### "two_column" — 两栏对比页
- 用于对比两个概念、展示优缺点、并列展示两类内容
- 必须包含：left_title、left_bullets、right_title、right_bullets
- 示例：对比"编译型语言"vs"解释型语言"

### "chart" — 图表页
- 用于展示数据、统计结果、趋势分析等
- 必须包含 chart 字段（见下方图表数据格式）
- 可同时包含 bullets 对图表进行说明
- 图表类型选择：bar（柱状图，适合对比）、pie（饼图，适合占比）、line（折线图，适合趋势）

### "summary" — 总结页
- 用于章节末尾的总结回顾
- 要点应简洁，突出核心结论
- 可包含"思考题"或"下一步"引导

## 每张幻灯片的详细要求

### title（标题）
- 简洁有力，通常 4-15 字
- 准确概括本页核心主题
- 各页标题之间逻辑递进

### bullets（要点列表）
- 每个要点必须是一句**完整的、有信息量的话**（15-50字），不能只是短语或标签
- 要点之间逻辑递进，由浅入深
- 要点应包含：概念定义、核心原理、公式或规则、实际例证、注意事项等
- 每页 3-5 个要点
- 严禁要点内容空洞，如"本节将介绍XX"这种无实质内容的句子
- 对于重要知识点，要点中应包含具体数值、公式、案例等可操作信息

### notes（讲解备注，用于语音旁白生成）
- 是一段**流畅完整的讲稿**（150-300字），适合真人朗读
- 对 bullets 中的每个要点进行深入展开和讲解
- 使用"我们""大家""同学们"等亲切称呼，增强课堂代入感
- 自然融入"例如""需要注意的是""思考一下"等教学引导语
- 适当加入设问、停顿提示（如"这里请大家思考一下"）
- 语言风格：专业但不生硬，亲切但不随意

## 图表数据格式（仅 layout="chart" 时需要）

当 layout="chart" 时，必须提供 chart 字段：

```json
"chart": {
  "type": "bar",           // 必填：bar（柱状图）、pie（饼图）、line（折线图）
  "title": "图表标题",      // 必填
  "categories": ["A", "B", "C"],  // 必填：X轴类别（pie图也需要）
  "series": [               // 必填：数据系列
    {"name": "系列1", "values": [10, 20, 30]}
  ]
}
```

## 两栏布局格式（仅 layout="two_column" 时需要）

```json
{
  "layout": "two_column",
  "title": "对比分析：XXX",
  "left_title": "概念A",
  "left_bullets": ["要点1", "要点2", "要点3"],
  "right_title": "概念B",
  "right_bullets": ["要点1", "要点2", "要点3"],
  "notes": "讲稿..."
}
```

## 知识点覆盖要求

1. 必须覆盖上下文中提供的**所有知识点**
2. 重要知识点（importance ≥ 0.8）应分配独立幻灯片或更多篇幅
3. 不同知识点之间应有自然过渡和逻辑衔接
4. 如果知识点之间存在前驱后继关系，应体现教学顺序

## 输出格式

{
  "title": "课件标题（通常是章节标题或更生动的变体）",
  "slides": [
    {
      "layout": "title",
      "title": "幻灯片标题",
      "notes": "讲稿..."
    },
    {
      "layout": "content",
      "title": "概念讲解",
      "bullets": ["包含具体内容的完整句子1", "包含具体内容的完整句子2"],
      "notes": "讲稿..."
    },
    {
      "layout": "two_column",
      "title": "对比分析",
      "left_title": "概念A",
      "left_bullets": ["要点1", "要点2"],
      "right_title": "概念B",
      "right_bullets": ["要点1", "要点2"],
      "notes": "讲稿..."
    },
    {
      "layout": "chart",
      "title": "数据分析",
      "chart": {
        "type": "bar",
        "title": "成绩分布",
        "categories": ["A", "B", "C"],
        "series": [{"name": "分数", "values": [90, 85, 92]}]
      },
      "bullets": ["结论1", "结论2"],
      "notes": "讲稿..."
    },
    {
      "layout": "summary",
      "title": "本章总结",
      "bullets": ["总结要点1", "总结要点2"],
      "notes": "讲稿..."
    }
  ]
}

**注意**：layout 字段是可选的，如果不提供则默认为 "content"。"""

_NARRATION_VOICE = "zh-CN-YunxiNeural"  # 旁白配音 - 男声


class PPTGenerator(BaseModalGenerator):
    modal_type = "ppt"

    async def generate(self, chapter_id: int, context: dict[str, Any]) -> dict[str, Any]:
        chapter_title = context.get("chapter_title", "未命名章节")
        kps = context.get("knowledge_points", [])

        # Build detailed knowledge point info for richer context
        kp_lines: list[str] = []
        for i, kp in enumerate(kps, 1):
            name = kp.get("name", "")
            kp_type = kp.get("type", "")
            importance = kp.get("importance", "")
            prerequisites = kp.get("prerequisites", [])
            extra = []
            if kp_type:
                extra.append(f"类型：{kp_type}")
            if importance != "":
                extra.append(f"重要性：{importance}")
            if prerequisites:
                extra.append(f"前置知识：{', '.join(prerequisites)}")
            suffix = f"（{'；'.join(extra)}）" if extra else ""
            kp_lines.append(f"{i}. {name}{suffix}")
        kp_text = "\n".join(kp_lines) if kp_lines else "无特定知识点"

        course_desc = context.get("course_description", "")
        desc_line = f"\n课程简介：{course_desc}" if course_desc else ""

        user = f"""## 课程信息
课程标题：{context.get("course_title", "")}{desc_line}

## 本章信息
章节标题：{chapter_title}

## 本章知识点（请全部覆盖）
{kp_text}

## 任务要求
请根据以上信息，生成一套教学PPT课件。要求：
1. 覆盖所有知识点，重要知识点给予更多篇幅
2. 幻灯片之间逻辑递进，有明确的教学主线
3. 每个要点都要有实质内容，不要空洞的占位文字
4. notes 备注要写成可直接朗读的讲稿"""

        try:
            raw = await self._call_llm(_PPT_SYSTEM, user)
            content = self._extract_json(raw)
        except Exception as e:
            logger.warning("[ppt] LLM generation failed: %s, using fallback", e)
            content = self._fallback_content(chapter_title, kps)

        course_id = context.get("course_id")

        # Run PPTX export and narration TTS in parallel
        pptx_task = asyncio.create_task(
            self._export_pptx_async(content, course_id, chapter_id, chapter_title)
        )
        narration_task = asyncio.create_task(
            self._synthesize_narrations(content, course_id, chapter_id, chapter_title)
        )

        file_path, narration_urls = await asyncio.gather(
            pptx_task, narration_task, return_exceptions=True,
        )

        if isinstance(file_path, Exception):
            logger.warning("[ppt] PPTX export failed: %s", file_path)
            file_path = None
        if isinstance(narration_urls, Exception):
            logger.warning("[ppt] Narration TTS failed: %s", narration_urls)
            narration_urls = []

        # Inject narration URLs into content JSON for frontend
        content["narration_urls"] = narration_urls

        return {
            "modal_type": self.modal_type,
            "content_json": json.dumps(content, ensure_ascii=False),
            "file_path": file_path,
        }

    # ------------------------------------------------------------------
    # Async wrapper for PPTX export (runs in thread pool)
    # ------------------------------------------------------------------
    async def _export_pptx_async(
        self,
        content: dict[str, Any],
        course_id: int | None,
        chapter_id: int,
        chapter_title: str,
    ) -> str | None:
        """Run PPTX export in a thread to avoid blocking the event loop."""
        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(
            None, self._export_pptx, content, course_id, chapter_id, chapter_title,
        )

    # ------------------------------------------------------------------
    # PPTX export — Academic Education style
    # ------------------------------------------------------------------
    def _export_pptx(
        self,
        content: dict[str, Any],
        course_id: int | None,
        chapter_id: int,
        chapter_title: str,
    ) -> str | None:
        """Export slides JSON to a beautifully designed .pptx file."""
        slides = content.get("slides", [])
        if not slides:
            return None

        prs = Presentation()
        prs.slide_width = _SLIDE_W
        prs.slide_height = _SLIDE_H

        total = len(slides)
        for idx, slide_data in enumerate(slides):
            layout = slide_data.get("layout", "content")
            page_num = idx + 1

            # First slide always uses title layout (unless explicit layout set)
            if idx == 0 and layout == "content":
                layout = "title"

            if layout == "title":
                self._add_title_slide_v2(prs, slide_data, chapter_title)
            elif layout == "content":
                self._add_content_slide_v2(prs, slide_data, chapter_title, page_num, total)
            elif layout == "two_column":
                self._add_two_column_slide(prs, slide_data, chapter_title, page_num, total)
            elif layout == "chart":
                self._add_chart_slide(prs, slide_data, chapter_title, page_num, total)
            elif layout == "summary":
                self._add_summary_slide(prs, slide_data, chapter_title, page_num, total)
            else:
                # Fallback to content layout
                self._add_content_slide_v2(prs, slide_data, chapter_title, page_num, total)

        # Save to assets directory
        course_dir_name = f"course_{course_id}" if course_id else f"chapter_{chapter_id}"
        export_dir = get_assets_root() / course_dir_name / "ppt"
        export_dir.mkdir(parents=True, exist_ok=True)

        safe_title = "".join(c for c in chapter_title if c.isalnum() or c in " _-")[:40].strip()
        filename = f"{safe_title or 'chapter'}_{chapter_id}.pptx"
        filepath = export_dir / filename

        prs.save(str(filepath))
        logger.info("[ppt] Exported .pptx to %s", filepath)

        return f"/assets/{course_dir_name}/ppt/{filename}"

    # ------------------------------------------------------------------
    # Title slide — gradient blue background, centered white title
    # ------------------------------------------------------------------
    def _add_title_slide(
        self,
        prs: Any,
        slide_data: dict[str, Any],
        chapter_title: str,
    ) -> None:
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Gradient background via stacked rectangles
        self._draw_gradient_bg(slide)

        # Decorative circle top-right
        self._add_decor_circle(slide, Inches(9.5), Inches(-1.5), Inches(4.0))

        # Decorative circle bottom-left
        self._add_decor_circle(slide, Inches(-1.5), Inches(5.5), Inches(3.0))

        # Centered title
        title_text = slide_data.get("title", chapter_title)
        left = Inches(1.5)
        top = Inches(2.2)
        width = Inches(10.3)
        height = Inches(1.8)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title_text
        self._set_font(run, Pt(40), _COLOR_WHITE, bold=True)

        # Subtitle / chapter label
        sub_left = Inches(1.5)
        sub_top = Inches(4.2)
        sub_width = Inches(10.3)
        sub_height = Inches(1.2)

        subBox = slide.shapes.add_textbox(sub_left, sub_top, sub_width, sub_height)
        stf = subBox.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        srun = sp.add_run()
        srun.text = chapter_title
        self._set_font(srun, Pt(22), _COLOR_ACCENT_YELLOW, bold=False)

        # Footer line
        self._add_footer_line(slide, chapter_title)

    # ------------------------------------------------------------------
    # Title slide V2 — illustration style
    # ------------------------------------------------------------------
    def _add_title_slide_v2(
        self,
        prs: Any,
        slide_data: dict[str, Any],
        chapter_title: str,
    ) -> None:
        """Illustration-style title slide with warm gradient, decorative blobs, and wave."""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Warm gradient background (orange -> light orange -> yellow)
        self._draw_gradient_bg_v2(slide)

        # Large decorative blobs (illustration style)
        # Top-right large circle
        self._add_decor_blob(slide, 8.5, -1.0, 5.0, 5.0, _COLOR_PRIMARY_LIGHT)
        # Bottom-left medium circle
        self._add_decor_blob(slide, -1.0, 4.5, 3.5, 3.5, _COLOR_ACCENT_YELLOW)
        # Center-right small circle
        self._add_decor_blob(slide, 7.0, 3.0, 2.0, 2.0, _COLOR_SECONDARY_LIGHT)

        # Title card (rounded rectangle with warm orange)
        card_left = Inches(1.0)
        card_top = Inches(1.8)
        card_width = Inches(11.3)
        card_height = Inches(2.5)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            card_left, card_top, card_width, card_height,
        )
        card.fill.solid()
        card.fill.fore_color.rgb = _COLOR_PRIMARY
        card.line.fill.background()
        # Rounded corners
        self._set_shape_radius(card, 0.15)

        # Title text (white, centered, large)
        title_text = slide_data.get("title", chapter_title)
        txBox = slide.shapes.add_textbox(
            card_left + Inches(0.3), card_top + Inches(0.3),
            card_width - Inches(0.6), card_height - Inches(0.6),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title_text
        self._set_font(run, Pt(38), _COLOR_WHITE, bold=True)

        # Subtitle / chapter label (below card)
        subBox = slide.shapes.add_textbox(
            Inches(1.0), Inches(4.8), Inches(11.3), Inches(0.8),
        )
        stf = subBox.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        srun = sp.add_run()
        srun.text = chapter_title
        self._set_font(srun, Pt(20), _COLOR_SECONDARY_DARK, bold=False)

        # Hand-drawn style wave at bottom
        self._add_decor_wave(slide, 6.6, _COLOR_ACCENT_YELLOW)

        # Footer line
        self._add_footer_line(slide, chapter_title)

    # ------------------------------------------------------------------
    # Content slide V2 — illustration style
    # ------------------------------------------------------------------
    def _add_content_slide_v2(
        self,
        prs: Any,
        slide_data: dict[str, Any],
        chapter_title: str,
        page_num: int,
        total: int,
    ) -> None:
        """Illustration-style content slide with warm bg, colored header, and icon badges."""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Warm white background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _COLOR_BG

        # Top header bar (gradient orange)
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), _SLIDE_W, Inches(1.2),
        )
        header.fill.solid()
        header.fill.fore_color.rgb = _COLOR_PRIMARY
        header.line.fill.background()

        # Header title text (white, left-aligned)
        h_title = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.15), Inches(10.5), Inches(0.9),
        )
        htf = h_title.text_frame
        htf.word_wrap = True
        hp = htf.paragraphs[0]
        hp.alignment = PP_ALIGN.LEFT
        hrun = hp.add_run()
        hrun.text = slide_data.get("title", "")
        self._set_font(hrun, Pt(28), _COLOR_WHITE, bold=True)

        # Left side decorative stripe (fresh green)
        self._add_side_decor_stripes(slide, side="left")

        # Right side decorative dots
        for i in range(3):
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(12.5), Inches(1.5 + i * 0.4),
                Inches(0.15), Inches(0.15),
            )
            dot.fill.solid()
            colors = [_COLOR_ACCENT_YELLOW, _COLOR_SECONDARY, _COLOR_ACCENT_PINK]
            dot.fill.fore_color.rgb = colors[i % len(colors)]
            dot.line.fill.background()

        # Bullet content area with numbered badges
        bullets = slide_data.get("bullets", [])
        content_left = Inches(1.5)
        content_top = Inches(1.8)
        content_width = Inches(10.5)

        # Light card background behind bullets
        if bullets:
            card_h = min(4.2, len(bullets) * 0.7 + 0.3)
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, content_left - Inches(0.2),
                content_top - Inches(0.15),
                content_width + Inches(0.4), Inches(card_h),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = _COLOR_WHITE
            card.line.color.rgb = _COLOR_SECONDARY_LIGHT
            card.line.width = Pt(1)
            self._set_shape_radius(card, 0.08)

        txBox = slide.shapes.add_textbox(
            Inches(2.0), Inches(2.0), Inches(9.8), Inches(4.0),
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, bullet_text in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.space_before = Pt(10)
            p.space_after = Pt(5)

            # Numbered badge (orange circle with number)
            self._add_number_badge(slide, 1.5, 2.0 + i * 0.65, i + 1, _COLOR_PRIMARY)

            # Bullet text
            run_text = p.add_run()
            run_text.text = bullet_text
            self._set_font(run_text, Pt(18), _COLOR_TEXT, bold=False)

        # Add a small decorative icon in bottom-right corner
        self._add_icon_shape(slide, 11.5, 6.2, 0.8, "bulb")

        # Speaker notes
        notes_text = slide_data.get("notes", "")
        if notes_text:
            try:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes_text
            except Exception:
                pass

        # Footer
        self._add_footer(slide, chapter_title, page_num, total)

    # ------------------------------------------------------------------
    # Two-column slide — illustration style
    # ------------------------------------------------------------------
    def _add_two_column_slide(
        self,
        prs: Any,
        slide_data: dict[str, Any],
        chapter_title: str,
        page_num: int,
        total: int,
    ) -> None:
        """Two-column layout for comparisons or parallel concepts."""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Warm white background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _COLOR_BG

        # Top header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), _SLIDE_W, Inches(1.1),
        )
        header.fill.solid()
        header.fill.fore_color.rgb = _COLOR_PRIMARY
        header.line.fill.background()

        # Header title
        h_title = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.15), Inches(10.5), Inches(0.8),
        )
        htf = h_title.text_frame
        hp = htf.paragraphs[0]
        hp.alignment = PP_ALIGN.LEFT
        hrun = hp.add_run()
        hrun.text = slide_data.get("title", "")
        self._set_font(hrun, Pt(26), _COLOR_WHITE, bold=True)

        # Left column background card (orange-tinted)
        left_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.4), Inches(1.3), Inches(6.0), Inches(5.2),
        )
        left_card.fill.solid()
        left_card.fill.fore_color.rgb = RGBColor(0xFF, 0xF5, 0xEB)  # Very light orange
        left_card.line.color.rgb = _COLOR_PRIMARY_LIGHT
        left_card.line.width = Pt(1)
        self._set_shape_radius(left_card, 0.08)

        # Right column background card (green-tinted)
        right_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.9), Inches(1.3), Inches(6.0), Inches(5.2),
        )
        right_card.fill.solid()
        right_card.fill.fore_color.rgb = RGBColor(0xE8, 0xF9, 0xF3)  # Very light green
        right_card.line.color.rgb = _COLOR_SECONDARY_LIGHT
        right_card.line.width = Pt(1)
        self._set_shape_radius(right_card, 0.08)

        # Left column title bar
        left_title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.4), Inches(1.3), Inches(6.0), Inches(0.6),
        )
        left_title_bar.fill.solid()
        left_title_bar.fill.fore_color.rgb = _COLOR_PRIMARY
        left_title_bar.line.fill.background()

        left_title_tx = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.35), Inches(5.6), Inches(0.5),
        )
        ltf = left_title_tx.text_frame
        lp = ltf.paragraphs[0]
        lp.alignment = PP_ALIGN.LEFT
        lr = lp.add_run()
        lr.text = slide_data.get("left_title", "左侧标题")
        self._set_font(lr, Pt(18), _COLOR_WHITE, bold=True)

        # Right column title bar
        right_title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.9), Inches(1.3), Inches(6.0), Inches(0.6),
        )
        right_title_bar.fill.solid()
        right_title_bar.fill.fore_color.rgb = _COLOR_SECONDARY
        right_title_bar.line.fill.background()

        right_title_tx = slide.shapes.add_textbox(
            Inches(7.1), Inches(1.35), Inches(5.6), Inches(0.5),
        )
        rtf = right_title_tx.text_frame
        rp = rtf.paragraphs[0]
        rp.alignment = PP_ALIGN.LEFT
        rr = rp.add_run()
        rr.text = slide_data.get("right_title", "右侧标题")
        self._set_font(rr, Pt(18), _COLOR_WHITE, bold=True)

        # Left column bullets
        left_bullets = slide_data.get("left_bullets", [])
        left_tx = slide.shapes.add_textbox(
            Inches(0.7), Inches(2.1), Inches(5.4), Inches(4.0),
        )
        ltf2 = left_tx.text_frame
        ltf2.word_wrap = True

        for i, bullet_text in enumerate(left_bullets):
            if i == 0:
                p = ltf2.paragraphs[0]
            else:
                p = ltf2.add_paragraph()
            p.space_before = Pt(8)
            p.space_after = Pt(4)
            run_b = p.add_run()
            run_b.text = "● "
            self._set_font(run_b, Pt(12), _COLOR_PRIMARY, bold=False)
            run_t = p.add_run()
            run_t.text = bullet_text
            self._set_font(run_t, Pt(16), _COLOR_TEXT, bold=False)

        # Right column bullets
        right_bullets = slide_data.get("right_bullets", [])
        right_tx = slide.shapes.add_textbox(
            Inches(7.1), Inches(2.1), Inches(5.4), Inches(4.0),
        )
        rtf2 = right_tx.text_frame
        rtf2.word_wrap = True

        for i, bullet_text in enumerate(right_bullets):
            if i == 0:
                p = rtf2.paragraphs[0]
            else:
                p = rtf2.add_paragraph()
            p.space_before = Pt(8)
            p.space_after = Pt(4)
            run_b = p.add_run()
            run_b.text = "● "
            self._set_font(run_b, Pt(12), _COLOR_SECONDARY, bold=False)
            run_t = p.add_run()
            run_t.text = bullet_text
            self._set_font(run_t, Pt(16), _COLOR_TEXT, bold=False)

        # Center divider decoration
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.65), Inches(1.3), Inches(0.08), Inches(5.2),
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = _COLOR_ACCENT_YELLOW
        divider.line.fill.background()

        # Speaker notes
        notes_text = slide_data.get("notes", "")
        if notes_text:
            try:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes_text
            except Exception:
                pass

        # Footer
        self._add_footer(slide, chapter_title, page_num, total)

    # ------------------------------------------------------------------
    # Chart slide — illustration style
    # ------------------------------------------------------------------
    def _add_chart_slide(
        self,
        prs: Any,
        slide_data: dict[str, Any],
        chapter_title: str,
        page_num: int,
        total: int,
    ) -> None:
        """Chart layout: left chart + right explanation."""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Warm white background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _COLOR_BG

        # Top header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), _SLIDE_W, Inches(1.1),
        )
        header.fill.solid()
        header.fill.fore_color.rgb = _COLOR_PRIMARY
        header.line.fill.background()

        # Header title
        h_title = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.15), Inches(10.5), Inches(0.8),
        )
        htf = h_title.text_frame
        hp = htf.paragraphs[0]
        hp.alignment = PP_ALIGN.LEFT
        hrun = hp.add_run()
        hrun.text = slide_data.get("title", "")
        self._set_font(hrun, Pt(26), _COLOR_WHITE, bold=True)

        # Chart area (left side)
        chart_data = slide_data.get("chart", {})
        if chart_data:
            self._add_chart_to_slide(slide, chart_data, 0.4, 1.3, 6.5, 5.2)
        else:
            # Placeholder if no chart data
            chart_ph = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.4), Inches(1.3), Inches(6.5), Inches(5.2),
            )
            chart_ph.fill.solid()
            chart_ph.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
            chart_ph.line.color.rgb = _COLOR_SECONDARY_LIGHT
            chart_ph.line.width = Pt(1)

        # Right side: explanation bullets
        bullets = slide_data.get("bullets", [])
        expl_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7.2), Inches(1.3), Inches(5.8), Inches(5.2),
        )
        expl_card.fill.solid()
        expl_card.fill.fore_color.rgb = _COLOR_WHITE
        expl_card.line.color.rgb = _COLOR_PRIMARY_LIGHT
        expl_card.line.width = Pt(1)
        self._set_shape_radius(expl_card, 0.08)

        expl_tx = slide.shapes.add_textbox(
            Inches(7.6), Inches(1.6), Inches(5.0), Inches(4.5),
        )
        etf = expl_tx.text_frame
        etf.word_wrap = True

        # Title for explanation area
        expl_title = slide.shapes.add_textbox(
            Inches(7.4), Inches(1.35), Inches(5.4), Inches(0.5),
        )
        ettf = expl_title.text_frame
        etp = ettf.paragraphs[0]
        etp.alignment = PP_ALIGN.LEFT
        etr = etp.add_run()
        etr.text = "📊 数据分析"
        self._set_font(etr, Pt(18), _COLOR_PRIMARY, bold=True)

        for i, bullet_text in enumerate(bullets):
            if i == 0:
                p = etf.paragraphs[0]
            else:
                p = etf.add_paragraph()
            p.space_before = Pt(10)
            p.space_after = Pt(5)
            run_b = p.add_run()
            run_b.text = "→ "
            self._set_font(run_b, Pt(14), _COLOR_SECONDARY, bold=False)
            run_t = p.add_run()
            run_t.text = bullet_text
            self._set_font(run_t, Pt(17), _COLOR_TEXT, bold=False)

        # Speaker notes
        notes_text = slide_data.get("notes", "")
        if notes_text:
            try:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes_text
            except Exception:
                pass

        # Footer
        self._add_footer(slide, chapter_title, page_num, total)

    # ------------------------------------------------------------------
    # Summary slide — illustration style
    # ------------------------------------------------------------------
    def _add_summary_slide(
        self,
        prs: Any,
        slide_data: dict[str, Any],
        chapter_title: str,
        page_num: int,
        total: int,
    ) -> None:
        """Summary layout with centered key points and icon badges."""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Warm gradient background
        self._draw_gradient_bg_v2(slide)

        # Translucent overlay for readability
        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), _SLIDE_W, _SLIDE_H,
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xF0)
        overlay.line.fill.background()
        # Send to back
        try:
            sp_tree = slide.shapes._spTree
            overlay_sp = overlay._element
            sp_tree.remove(overlay_sp)
            sp_tree.insert(1, overlay_sp)
        except Exception:
            pass

        # Title area (centered, orange background)
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2.5), Inches(0.4), Inches(8.3), Inches(1.0),
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = _COLOR_PRIMARY
        title_bg.line.fill.background()
        self._set_shape_radius(title_bg, 0.1)

        title_tx = slide.shapes.add_textbox(
            Inches(2.5), Inches(0.45), Inches(8.3), Inches(0.9),
        )
        ttf = title_tx.text_frame
        tp = ttf.paragraphs[0]
        tp.alignment = PP_ALIGN.CENTER
        tr = tp.add_run()
        tr.text = slide_data.get("title", "本章总结")
        self._set_font(tr, Pt(30), _COLOR_WHITE, bold=True)

        # Summary bullets with icon badges
        bullets = slide_data.get("bullets", [])
        num_bullets = len(bullets)

        # Calculate positions for bullet cards
        card_h = min(0.9, 4.5 / max(num_bullets, 1))
        start_y = 1.6

        icon_types = ["bulb", "check", "star", "gear", "heart"]

        for i, bullet_text in enumerate(bullets):
            y = start_y + i * (card_h + 0.15)

            # Bullet card
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(2.0), Inches(y), Inches(9.3), Inches(card_h),
            )
            card.fill.solid()
            # Alternate colors
            colors = [_COLOR_PRIMARY_LIGHT, _COLOR_SECONDARY_LIGHT, _COLOR_ACCENT_YELLOW]
            card.fill.fore_color.rgb = colors[i % 3]
            card.line.fill.background()
            self._set_shape_radius(card, 0.06)

            # Icon badge
            self._add_icon_shape(slide, 1.2, float(y) + 0.05, 0.6, icon_types[i % len(icon_types)])

            # Bullet text
            txBox = slide.shapes.add_textbox(
                Inches(2.3), Inches(y + 0.1), Inches(8.8), Inches(card_h - 0.2),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = bullet_text
            self._set_font(run, Pt(18), _COLOR_TEXT, bold=False)

        # Decorative wave at bottom
        self._add_decor_wave(slide, 6.7, _COLOR_ACCENT_YELLOW)

        # Speaker notes
        notes_text = slide_data.get("notes", "")
        if notes_text:
            try:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes_text
            except Exception:
                pass

        # Footer
        self._add_footer(slide, chapter_title, page_num, total)

    # ------------------------------------------------------------------
    # Content slide — header bar + left decor + content card (legacy)
    # ------------------------------------------------------------------
    def _add_content_slide(
        self,
        prs: Any,
        slide_data: dict[str, Any],
        chapter_title: str,
        page_num: int,
        total: int,
    ) -> None:
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Warm white background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _COLOR_BG

        # Top header bar (warm orange gradient feel)
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), _SLIDE_W, Inches(1.15),
        )
        header.fill.solid()
        header.fill.fore_color.rgb = _COLOR_PRIMARY
        header.line.fill.background()  # no border

        # Header title text
        h_title = slide.shapes.add_textbox(
            Inches(1.0), Inches(0.15), Inches(10.5), Inches(0.85),
        )
        htf = h_title.text_frame
        htf.word_wrap = True
        hp = htf.paragraphs[0]
        hp.alignment = PP_ALIGN.LEFT
        hrun = hp.add_run()
        hrun.text = slide_data.get("title", "")
        self._set_font(hrun, Pt(30), _COLOR_WHITE, bold=True)

        # Left decorative bar (fresh green)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.15), Inches(0.12), Inches(5.85),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = _COLOR_SECONDARY
        bar.line.fill.background()

        # Bottom accent line (thin, warm yellow)
        bot_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(7.0), _SLIDE_W, Inches(0.04),
        )
        bot_line.fill.solid()
        bot_line.fill.fore_color.rgb = _COLOR_ACCENT_YELLOW
        bot_line.line.fill.background()

        # Bullet content area
        bullets = slide_data.get("bullets", [])
        content_left = Inches(1.2)
        content_top = Inches(1.7)
        content_width = Inches(10.8)

        # Light card background behind bullets
        if bullets:
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, content_left, content_top,
                content_width, Inches(min(4.5, len(bullets) * 0.65 + 0.5)),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = _COLOR_WHITE
            card.line.color.rgb = _COLOR_SECONDARY_LIGHT
            card.line.width = Pt(1)
            # Rounded corners
            self._set_shape_radius(card)

        txBox = slide.shapes.add_textbox(
            Inches(1.8), Inches(2.0), Inches(9.8), Inches(4.2),
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, bullet_text in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.space_before = Pt(12)
            p.space_after = Pt(6)

            # Custom bullet dot
            run_bullet = p.add_run()
            run_bullet.text = "● "
            self._set_font(run_bullet, Pt(14), _COLOR_BULLET, bold=False)

            # Bullet text
            run_text = p.add_run()
            run_text.text = bullet_text
            self._set_font(run_text, Pt(20), _COLOR_TEXT, bold=False)

        # Speaker notes
        notes_text = slide_data.get("notes", "")
        if notes_text:
            try:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes_text  # type: ignore[union-attr]
            except Exception:
                pass

        # Footer
        self._add_footer(slide, chapter_title, page_num, total)

    # ------------------------------------------------------------------
    # Illustration-style decoration helpers
    # ------------------------------------------------------------------
    def _add_decor_blob(self, slide: Any, left: float, top: float,
                        width: float, height: float, color: Any,
                        _opacity: float = 1.0) -> None:
        """Add a soft blob shape (illustration style decorative element)."""
        blob = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(left), Inches(top),
            Inches(width), Inches(height),
        )
        blob.fill.solid()
        blob.fill.fore_color.rgb = color
        blob.line.fill.background()
        # Make it slightly transparent by using a lighter version of the color
        # (python-pptx doesn't support alpha, so we simulate with lighter color)

    def _add_decor_wave(self, slide: Any, y_pos: float,
                         color: Any = None) -> None:
        """Add a hand-drawn style wave line at the bottom of title slides."""
        if color is None:
            color = _COLOR_ACCENT_YELLOW
        # Draw a wavy line using multiple small arcs (simulated with rectangle rows)
        wave_h = Inches(0.08)
        wave_w = Inches(0.5)
        y = Inches(y_pos)
        for i in range(int(_SLIDE_W / wave_w) + 1):
            x = i * wave_w
            offset = Inches(0.03) if (i % 2 == 0) else Inches(0)
            wave_piece = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, y + offset, wave_w, wave_h,
            )
            wave_piece.fill.solid()
            wave_piece.fill.fore_color.rgb = color
            wave_piece.line.fill.background()
            # Slightly randomize corner radius for hand-drawn feel
            try:
                spPr = wave_piece._element.find(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}spPr"
                )
                if spPr is not None:
                    prstGeom = spPr.find(
                        "{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom"
                    )
                    if prstGeom is not None:
                        prstGeom.set("prst", "roundRect")
            except Exception:
                pass

    def _add_icon_shape(self, slide: Any, left: float, top: float,
                         size: float, icon_type: str) -> None:
        """Add a simple icon shape (illustration style).
        
        icon_type: "bulb", "gear", "star", "heart", "arrow", "check"
        """
        c = RGBColor(0xFF, 0xFF, 0xFF)  # white icon
        bg = RGBColor(0xFF, 0x8C, 0x42)  # orange background circle
        
        # Background circle
        bg_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size),
        )
        bg_circle.fill.solid()
        bg_circle.fill.fore_color.rgb = bg
        bg_circle.line.fill.background()
        
        # Icon label (use text as simple icon replacement)
        icon_labels = {
            "bulb": "💡",
            "gear": "⚙",
            "star": "★",
            "heart": "♥",
            "arrow": "→",
            "check": "✓",
            "book": "📖",
            "idea": "💡",
            "code": "{ }",
            "data": "📊",
        }
        label = icon_labels.get(icon_type, "●")
        
        txBox = slide.shapes.add_textbox(
            Inches(left + size * 0.2), Inches(top + size * 0.2),
            Inches(size * 0.6), Inches(size * 0.6),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(int(size * 40))
        run.font.color.rgb = c
        run.font.bold = True

    def _add_number_badge(self, slide: Any, left: float, top: float,
                            number: int, color: Any = None) -> None:
        """Add a numbered badage circle (e.g., 1, 2, 3...) for bullet points."""
        if color is None:
            color = _COLOR_PRIMARY
        
        # Circle background
        size = 0.45
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        
        # Number text
        txBox = slide.shapes.add_textbox(
            Inches(left + 0.05), Inches(top + 0.05),
            Inches(size - 0.1), Inches(size - 0.1),
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(number)
        run.font.size = Pt(18)
        run.font.color.rgb = _COLOR_WHITE
        run.font.bold = True

    def _add_side_decor_stripes(self, slide: Any, side: str = "left") -> None:
        """Add colorful decorative stripes on the side of the slide."""
        colors = [_COLOR_PRIMARY, _COLOR_SECONDARY, _COLOR_ACCENT_YELLOW, _COLOR_ACCENT_PINK]
        stripe_w = Inches(0.06)
        if side == "left":
            x = Inches(0)
        else:
            x = _SLIDE_W - stripe_w
        
        for i, color in enumerate(colors):
            y = Inches(1.5 + i * 0.35)
            stripe = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, stripe_w, Inches(0.25),
            )
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = color
            stripe.line.fill.background()

    # ------------------------------------------------------------------
    # Illustration-style drawing helpers
    # ------------------------------------------------------------------
    def _draw_gradient_bg_v2(self, slide: Any) -> None:
        """Draw a warm illustration-style gradient background.
        
        Uses overlapping translucent rectangles to create a smooth
        warm gradient from orange (top) to yellow (bottom).
        """
        steps = 40
        band_h = _SLIDE_H / steps

        for i in range(steps):
            t = i / (steps - 1)
            # Orange -> Yellow gradient
            r = int(0xFF + (0xFF - 0xFF) * t)
            g = int(0x8C + (0xE6 - 0x8C) * t)
            b = int(0x42 + (0x6D - 0x42) * t)

            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(i) * band_h, _SLIDE_W, band_h + Inches(0.01),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(r, g, b)
            shape.line.fill.background()

        # Add a soft light overlay at top for depth
        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), _SLIDE_W, Inches(2.0),
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0xFF, 0xBF, 0x80)  # Light orange
        overlay.line.fill.background()
        # Move to back
        try:
            sp_tree = slide.shapes._spTree
            overlay_sp = overlay._element
            sp_tree.remove(overlay_sp)
            sp_tree.insert(1, overlay_sp)
        except Exception:
            pass

    def _set_shape_radius(self, shape: Any, _radius: float = 0.1) -> None:
        """Set rounded corner radius on a shape (improved version)."""
        try:
            spPr = shape._element.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}spPr"
            )
            if spPr is not None:
                prstGeom = spPr.find(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom"
                )
                if prstGeom is not None:
                    prstGeom.set("prst", "roundRect")
                    # Set radius via avLst
                    avLst = prstGeom.find(
                        "{http://schemas.openxmlformats.org/drawingml/2006/main}avLst"
                    )
                    if avLst is None:
                        avLst = prstGeom.makeelement(
                            "{http://schemas.openxmlformats.org/drawingml/2006/main}avLst",
                            {},
                        )
                        prstGeom.append(avLst)
        except Exception:
            pass

    # ------------------------------------------------------------------
    # Drawing helpers (legacy)
    # ------------------------------------------------------------------
    def _draw_gradient_bg(self, slide: Any) -> None:
        """Draw a vertical gradient from warm orange (top) to light orange (bottom) using rectangles."""
        steps = 30
        band_h = int(_SLIDE_H / steps)

        for i in range(steps):
            t = i / (steps - 1)
            r = int(_COLOR_PRIMARY_DARK[0] + (_COLOR_PRIMARY[0] - _COLOR_PRIMARY_DARK[0]) * t)
            g = int(_COLOR_PRIMARY_DARK[1] + (_COLOR_PRIMARY[1] - _COLOR_PRIMARY_DARK[1]) * t)
            b = int(_COLOR_PRIMARY_DARK[2] + (_COLOR_PRIMARY[2] - _COLOR_PRIMARY_DARK[2]) * t)

            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), band_h * i, _SLIDE_W, band_h,
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(r, g, b)
            shape.line.fill.background()

    @staticmethod
    def _add_decor_circle(slide: Any, left: Any, top: Any, size: Any) -> None:
        """Add a semi-transparent decorative circle (illustration style)."""
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left, top, size, size,
        )
        circle.fill.solid()
        # Warm yellow semi-transparent illusion
        circle.fill.fore_color.rgb = _COLOR_DECOR_CIRCLE
        circle.line.fill.background()

    def _add_footer_line(self, slide: Any, chapter_title: str) -> None:
        """Add a thin footer line with chapter title on title slide."""
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.5), Inches(6.4), Inches(10.3), Inches(0.015),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = _COLOR_ACCENT_YELLOW
        line.line.fill.background()

        ft = slide.shapes.add_textbox(
            Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.5),
        )
        tf = ft.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = chapter_title
        self._set_font(run, Pt(12), _COLOR_ACCENT_YELLOW, bold=False)

    def _add_footer(self, slide: Any, chapter_title: str, page_num: int, total: int) -> None:
        """Add footer with chapter title and page number."""
        # Chapter title on left
        ft_left = slide.shapes.add_textbox(
            Inches(0.5), Inches(7.05), Inches(8.0), Inches(0.35),
        )
        tf = ft_left.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = chapter_title
        self._set_font(run, Pt(10), _COLOR_TEXT_LIGHT, bold=False)

        # Page number on right
        ft_right = slide.shapes.add_textbox(
            Inches(10.0), Inches(7.05), Inches(3.0), Inches(0.35),
        )
        tf2 = ft_right.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        run2 = p2.add_run()
        run2.text = f"{page_num} / {total}"
        self._set_font(run2, Pt(10), _COLOR_TEXT_LIGHT, bold=False)

    @staticmethod
    def _set_font(run: Any, size: Any, color: Any, *, bold: bool = False) -> None:
        """Apply font settings to a run."""
        run.font.size = size
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = _FONT_FAMILY
        # Set East Asian font
        rPr = run._r.get_or_add_rPr()
        ea = rPr.makeelement(qn("a:ea"), {})
        ea.set("typeface", _FONT_FAMILY)
        rPr.append(ea)


    # ------------------------------------------------------------------
    # Chart drawing helper
    # ------------------------------------------------------------------
    def _add_chart_to_slide(self, slide: Any, chart_data: dict,
                           left: float, top: float,
                           width: float, height: float) -> None:
        """Draw a chart on the slide using python-pptx chart API."""
        try:
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE

            cd = CategoryChartData()
            cd.categories = chart_data.get("categories", [])

            for series in chart_data.get("series", []):
                cd.add_series(series["name"], series["values"])

            type_map = {
                "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "pie": XL_CHART_TYPE.PIE,
                "line": XL_CHART_TYPE.LINE,
            }
            chart_type = type_map.get(
                chart_data.get("type", "bar"),
                XL_CHART_TYPE.COLUMN_CLUSTERED,
            )

            chart_shape = slide.shapes.add_chart(
                chart_type,
                Inches(left), Inches(top), Inches(width), Inches(height),
                cd,
            )

            # Apply illustration-style colors to chart
            try:
                chart = chart_shape.chart
                colors = [
                    _COLOR_PRIMARY,
                    _COLOR_SECONDARY,
                    _COLOR_ACCENT_YELLOW,
                    _COLOR_ACCENT_PINK,
                    _COLOR_SECONDARY_LIGHT,
                ]
                for i, series in enumerate(chart.series):
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = colors[i % len(colors)]
            except Exception:
                pass  # Chart styling is optional

        except ImportError:
            # python-pptx chart support not available
            pass

    # ------------------------------------------------------------------
    # Narration TTS synthesis
    # ------------------------------------------------------------------
    async def _synthesize_narrations(
        self,
        content: dict[str, Any],
        course_id: int | None,
        chapter_id: int,
        chapter_title: str,
    ) -> list[str]:
        """Synthesize narration MP3 for each slide via Edge TTS."""
        slides = content.get("slides", [])
        if not slides:
            return []

        tts = get_tts_provider()
        course_dir_name = f"course_{course_id}" if course_id else f"chapter_{chapter_id}"
        export_dir = get_assets_root() / course_dir_name / "narration"
        export_dir.mkdir(parents=True, exist_ok=True)

        safe_title = "".join(c for c in chapter_title if c.isalnum() or c in " _-")[:40].strip()
        urls: list[str] = []

        for i, slide in enumerate(slides):
            notes = slide.get("notes", "")
            if not notes.strip():
                urls.append("")
                continue

            try:
                audio_bytes = await tts.synthesize(notes, {"voice": _NARRATION_VOICE})
                filename = f"{safe_title or 'ch'}_{chapter_id}_slide_{i}.mp3"
                filepath = export_dir / filename
                filepath.write_bytes(audio_bytes)
                urls.append(f"/assets/{course_dir_name}/narration/{filename}")
                logger.debug("[ppt] Narration MP3 for slide %d: %s", i, filename)
            except Exception as e:
                logger.warning("[ppt] Narration TTS failed for slide %d: %s", i, e)
                urls.append("")

            # Small delay between TTS calls
            await asyncio.sleep(0.15)

        logger.info("[ppt] Generated %d narration MP3s", len([u for u in urls if u]))
        return urls

    def _fallback_content(self, chapter_title: str, kps: list[dict[str, Any]]) -> dict[str, Any]:
        kp_names = [kp.get("name", "") for kp in kps]
        kp_types = {kp.get("name", ""): kp.get("type", "概念") for kp in kps}

        slides: list[dict[str, Any]] = []

        # Slide 1: Chapter overview (title layout)
        slides.append({
            "layout": "title",
            "title": f"{chapter_title} - 概述",
            "notes": f"大家好！欢迎来到{chapter_title}的学习。本章我们将系统学习{', '.join(kp_names[:3])}等重要知识点。通过本章的学习，你将建立起完整的知识框架，并能将其应用到实际问题中。",
        })

        # Slide 2-N: Distribute knowledge points across slides (max 3 KPs per slide)
        kps_per_slide = 3
        for start in range(0, len(kp_names), kps_per_slide):
            chunk = kp_names[start:start + kps_per_slide]
            if not chunk:
                continue
            kp_items = []
            for name in chunk:
                kp_type = kp_types.get(name, "概念")
                kp_items.append(f"{name}（{kp_type}）：理解其核心定义、原理和应用场景")
            slides.append({
                "layout": "content",
                "title": "核心知识点" if start == 0 else "知识拓展",
                "bullets": kp_items,
                "notes": f"接下来我们重点学习{'、'.join(chunk)}。{'对于每个知识点，我们不仅要理解概念定义，更要掌握其背后的原理和实际应用方法。' if len(chunk) > 1 else '让我们深入探讨这个知识点的核心内容。'}请同学们注意做好笔记。",
            })

        # Slide last: Summary
        slides.append({
            "layout": "summary",
            "title": "本章总结",
            "bullets": [
                f"回顾本章核心：{', '.join(kp_names[:3])}{'等' if len(kp_names) > 3 else ''}",
                "梳理知识体系，建立知识点之间的联系",
                "完成章节练习，检验学习效果",
                "思考：如何将本章知识应用到实际问题中？",
            ],
            "notes": f"好的，以上就是{chapter_title}的全部内容。我们学习了{', '.join(kp_names[:3])}等核心知识点。建议同学们课后整理笔记，建立自己的知识框架，并通过练习题来检验掌握程度。有任何疑问欢迎随时提问，我们下次课再见！",
        })

        return {
            "title": chapter_title,
            "slides": slides,
        }
